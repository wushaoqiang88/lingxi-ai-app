"""Real document file generation: PPT, PDF, Excel/CSV."""

from __future__ import annotations

import json
import os
import re
import uuid
from pathlib import Path
from typing import Any

# ── PPT ──────────────────────────────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ── Excel ────────────────────────────────────────────────────────────
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

# ── PDF ──────────────────────────────────────────────────────────────
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import mm
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib.colors import HexColor

# ── output directory ─────────────────────────────────────────────────
OUTPUT_DIR = Path(__file__).resolve().parents[1] / "generated_docs"
OUTPUT_DIR.mkdir(exist_ok=True)


def _uid() -> str:
    return uuid.uuid4().hex[:12]


# =====================================================================
#  Detect document type from user text
# =====================================================================

def detect_doc_type(user_text: str) -> str | None:
    """Return 'pptx' | 'pdf' | 'xlsx' | None based on keywords in user input."""
    t = user_text.lower()
    if any(k in t for k in ("ppt", "演示", "幻灯片", "slide")):
        return "pptx"
    if any(k in t for k in ("pdf", "报告", "文档")):
        return "pdf"
    if any(k in t for k in ("表格", "excel", "xlsx", "csv", "数据表")):
        return "xlsx"
    return None


# =====================================================================
#  Build system prompt that asks AI to return structured JSON
# =====================================================================

DOC_JSON_PROMPTS = {
"pptx": """你是世界顶级PPT策划与视觉设计师。用户要求生成一份专业级PPT演示文稿。
请根据用户需求，输出 **严格 JSON**（不要 markdown 代码块，不要多余解释）。格式：
{
  "title": "演示文稿标题",
  "subtitle": "一句话副标题",
  "author": "灵犀 AI",
  "theme": "从以下主题中选一个最适合内容的：科技蓝/商务绿/活力橙/高级灰/极简白/深邃紫/中国红/渐变蓝紫",
  "cover_style": "从 geometric/wave/minimal 中选一个封面风格",
  "slides": [
    {
      "title": "幻灯片标题",
      "kicker": "短标签（如 市场洞察/商业模式/落地路径）",
      "visual": "从 cards/timeline/process/metrics/comparison/icon_grid 中选择最适合的",
      "bullets": ["要点1", "要点2", "要点3"],
      "note": "可选的补充说明文字，出现在页面底部"
    }
  ]
}
要求：
- 幻灯片不少于 7 页，每页 3-5 个要点
- 内容专业、结构清晰，使用适合视觉排版的精炼短句
- visual 类型要根据内容特点选择，不要全部用 cards
- 偶数位置的 slide 尽量选不同的 visual 类型以增加多样性""",

    "pdf": """你是职场文档助手。用户要求生成 PDF 文档。
请根据用户需求，输出 **严格 JSON**（不要 markdown 代码块，不要多余解释）。格式：
{
  "title": "文档标题",
  "sections": [
    {
      "heading": "章节标题",
      "paragraphs": ["段落1正文...", "段落2正文..."]
    }
  ]
}
不少于 4 个章节，每个章节 1-3 段正文。内容专业、信息密度高。""",

    "xlsx": """你是职场文档助手。用户要求生成 Excel 表格。
请根据用户需求，输出 **严格 JSON**（不要 markdown 代码块，不要多余解释）。格式：
{
  "title": "表格标题",
  "headers": ["列名1", "列名2", "列名3"],
  "rows": [
    ["数据1", "数据2", "数据3"],
    ["数据4", "数据5", "数据6"]
  ]
}
表头不少于 3 列，数据行不少于 5 行。内容真实合理。""",
}


def get_doc_system_prompt(doc_type: str) -> str:
    return DOC_JSON_PROMPTS[doc_type]


# =====================================================================
#  Parse AI response to JSON
# =====================================================================

def parse_ai_json(raw: str) -> dict[str, Any]:
    """Try to extract JSON from AI response text."""
    # Strip markdown code fences if any
    cleaned = re.sub(r"^```(?:json)?\s*", "", raw.strip())
    cleaned = re.sub(r"\s*```$", "", cleaned)
    return json.loads(cleaned)


# =====================================================================
#  Register CJK font for PDF
# =====================================================================

_font_registered = False


def _register_cjk_font() -> str:
    """Register a CJK font for reportlab. Return font name."""
    global _font_registered
    font_name = "CJKFont"
    if _font_registered:
        return font_name

    # Try common CJK font paths on macOS / Linux
    candidates = [
        "/System/Library/Fonts/PingFang.ttc",
        "/System/Library/Fonts/STHeiti Light.ttc",
        "/System/Library/Fonts/Hiragino Sans GB.ttc",
        "/System/Library/Fonts/Supplemental/Songti.ttc",
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
    ]
    for path in candidates:
        if os.path.exists(path):
            try:
                pdfmetrics.registerFont(TTFont(font_name, path))
                _font_registered = True
                return font_name
            except Exception:
                continue
    # Fallback to Helvetica (ASCII only, CJK will show boxes)
    return "Helvetica"


# =====================================================================
#  PPT generator — Professional Design System
# =====================================================================

# Slide dimensions (16:9 widescreen)
_SW = Inches(13.333)
_SH = Inches(7.5)


def generate_pptx(data: dict[str, Any]) -> Path:
    prs = Presentation()
    prs.slide_width = _SW
    prs.slide_height = _SH

    title_text = data.get("title", "演示文稿")
    subtitle_text = data.get("subtitle", "由灵犀 AI 自动生成")
    author = data.get("author", "灵犀 AI")
    theme = _ppt_theme(data.get("theme", "科技蓝"))
    cover_style = data.get("cover_style", "geometric")
    slides_data = data.get("slides", [])

    # ── Cover ──
    _add_cover(prs, title_text, subtitle_text, author, theme, cover_style)

    # ── Table of contents ──
    if len(slides_data) >= 4:
        _add_toc_slide(prs, slides_data, theme)

    # ── Content slides with section dividers ──
    total = len(slides_data)
    for index, item in enumerate(slides_data, start=1):
        # Add section divider every 3 slides
        if index > 1 and (index - 1) % 3 == 0:
            _add_section_divider(prs, item, index, total, theme)

        visual = item.get("visual", "cards")
        layout_map = {
            "cards": _add_cards_slide,
            "timeline": _add_timeline_slide,
            "metrics": _add_metrics_slide,
            "process": _add_process_slide,
            "comparison": _add_comparison_slide,
            "icon_grid": _add_icon_grid_slide,
        }
        layout_fn = layout_map.get(visual, _add_cards_slide)
        layout_fn(prs, item, index, total, theme)

    # ── Closing ──
    _add_closing(prs, title_text, author, theme)

    filename = f"{_uid()}_{_safe(title_text)}.pptx"
    path = OUTPUT_DIR / filename
    prs.save(str(path))
    return path


# ── Theme system (8 themes) ─────────────────────────────────────────

def _ppt_theme(name: str) -> dict[str, Any]:
    presets = {
        "科技蓝": {
            "bg": RGBColor(245, 248, 255), "bg_dark": RGBColor(15, 23, 42),
            "primary": RGBColor(37, 99, 235), "primary_light": RGBColor(191, 219, 254),
            "secondary": RGBColor(6, 182, 212), "accent": RGBColor(99, 102, 241),
            "dark": RGBColor(15, 23, 42), "muted": RGBColor(100, 116, 139),
            "card": RGBColor(255, 255, 255), "card_alt": RGBColor(241, 245, 249),
            "border": RGBColor(226, 232, 240), "text_light": RGBColor(203, 213, 225),
            "gradient": [RGBColor(37, 99, 235), RGBColor(99, 102, 241)],
        },
        "商务绿": {
            "bg": RGBColor(244, 253, 249), "bg_dark": RGBColor(6, 78, 59),
            "primary": RGBColor(5, 150, 105), "primary_light": RGBColor(167, 243, 208),
            "secondary": RGBColor(20, 184, 166), "accent": RGBColor(16, 185, 129),
            "dark": RGBColor(15, 23, 42), "muted": RGBColor(71, 85, 105),
            "card": RGBColor(255, 255, 255), "card_alt": RGBColor(236, 253, 245),
            "border": RGBColor(209, 250, 229), "text_light": RGBColor(167, 243, 208),
            "gradient": [RGBColor(5, 150, 105), RGBColor(20, 184, 166)],
        },
        "活力橙": {
            "bg": RGBColor(255, 247, 237), "bg_dark": RGBColor(124, 45, 18),
            "primary": RGBColor(234, 88, 12), "primary_light": RGBColor(254, 215, 170),
            "secondary": RGBColor(245, 158, 11), "accent": RGBColor(239, 68, 68),
            "dark": RGBColor(28, 25, 23), "muted": RGBColor(87, 83, 78),
            "card": RGBColor(255, 255, 255), "card_alt": RGBColor(255, 247, 237),
            "border": RGBColor(254, 215, 170), "text_light": RGBColor(254, 215, 170),
            "gradient": [RGBColor(234, 88, 12), RGBColor(245, 158, 11)],
        },
        "高级灰": {
            "bg": RGBColor(248, 250, 252), "bg_dark": RGBColor(30, 41, 59),
            "primary": RGBColor(51, 65, 85), "primary_light": RGBColor(203, 213, 225),
            "secondary": RGBColor(99, 102, 241), "accent": RGBColor(139, 92, 246),
            "dark": RGBColor(15, 23, 42), "muted": RGBColor(100, 116, 139),
            "card": RGBColor(255, 255, 255), "card_alt": RGBColor(241, 245, 249),
            "border": RGBColor(226, 232, 240), "text_light": RGBColor(148, 163, 184),
            "gradient": [RGBColor(51, 65, 85), RGBColor(99, 102, 241)],
        },
        "极简白": {
            "bg": RGBColor(255, 255, 255), "bg_dark": RGBColor(24, 24, 27),
            "primary": RGBColor(24, 24, 27), "primary_light": RGBColor(228, 228, 231),
            "secondary": RGBColor(113, 113, 122), "accent": RGBColor(37, 99, 235),
            "dark": RGBColor(24, 24, 27), "muted": RGBColor(161, 161, 170),
            "card": RGBColor(250, 250, 250), "card_alt": RGBColor(244, 244, 245),
            "border": RGBColor(228, 228, 231), "text_light": RGBColor(161, 161, 170),
            "gradient": [RGBColor(24, 24, 27), RGBColor(63, 63, 70)],
        },
        "深邃紫": {
            "bg": RGBColor(250, 245, 255), "bg_dark": RGBColor(59, 7, 100),
            "primary": RGBColor(139, 92, 246), "primary_light": RGBColor(221, 214, 254),
            "secondary": RGBColor(192, 132, 252), "accent": RGBColor(236, 72, 153),
            "dark": RGBColor(30, 10, 60), "muted": RGBColor(107, 114, 128),
            "card": RGBColor(255, 255, 255), "card_alt": RGBColor(245, 243, 255),
            "border": RGBColor(221, 214, 254), "text_light": RGBColor(196, 181, 253),
            "gradient": [RGBColor(139, 92, 246), RGBColor(236, 72, 153)],
        },
        "中国红": {
            "bg": RGBColor(255, 241, 242), "bg_dark": RGBColor(127, 29, 29),
            "primary": RGBColor(220, 38, 38), "primary_light": RGBColor(254, 202, 202),
            "secondary": RGBColor(239, 68, 68), "accent": RGBColor(245, 158, 11),
            "dark": RGBColor(69, 10, 10), "muted": RGBColor(107, 114, 128),
            "card": RGBColor(255, 255, 255), "card_alt": RGBColor(254, 242, 242),
            "border": RGBColor(254, 202, 202), "text_light": RGBColor(252, 165, 165),
            "gradient": [RGBColor(220, 38, 38), RGBColor(185, 28, 28)],
        },
        "渐变蓝紫": {
            "bg": RGBColor(238, 242, 255), "bg_dark": RGBColor(30, 27, 75),
            "primary": RGBColor(79, 70, 229), "primary_light": RGBColor(199, 210, 254),
            "secondary": RGBColor(168, 85, 247), "accent": RGBColor(6, 182, 212),
            "dark": RGBColor(30, 27, 75), "muted": RGBColor(100, 116, 139),
            "card": RGBColor(255, 255, 255), "card_alt": RGBColor(238, 242, 255),
            "border": RGBColor(199, 210, 254), "text_light": RGBColor(165, 180, 252),
            "gradient": [RGBColor(79, 70, 229), RGBColor(168, 85, 247)],
        },
    }
    return presets.get(name, presets["科技蓝"])


# ── Low-level drawing helpers ────────────────────────────────────────

def _blank(prs, theme, dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = theme["bg_dark"] if dark else theme["bg"]
    return slide


def _rect(slide, l, t, w, h, fill, line=None, radius=False):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(st, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp


def _circle(slide, l, t, d, fill, alpha_pct=100):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


def _diamond(slide, l, t, s, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(l), Inches(t), Inches(s), Inches(s))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


def _hexagon(slide, l, t, w, h, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


def _triangle(slide, l, t, w, h, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


def _arrow_right(slide, l, t, w, h, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


def _txt(slide, text, l, t, w, h, size, color, bold=False, align=None, italic=False):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    run = p.runs[0]
    run.font.name = "PingFang SC"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


# ── Decorative patterns ─────────────────────────────────────────────

def _decor_dots(slide, x, y, theme, count=6):
    """Add small decorative dots."""
    for i in range(count):
        _circle(slide, x + i * 0.35, y, 0.12, theme["primary_light"])


def _decor_line(slide, x, y, w, theme, thin=True):
    """Add a decorative horizontal line."""
    h = 0.035 if thin else 0.07
    _rect(slide, x, y, w, h, theme["primary"])


def _decor_corner_shapes(slide, theme):
    """Add decorative shapes in corners."""
    _circle(slide, 11.8, -0.5, 2.0, theme["primary_light"])
    _circle(slide, 12.3, 0.0, 0.7, theme["secondary"])
    _diamond(slide, -0.3, 6.2, 1.0, theme["primary_light"])


def _progress_bar(slide, index, total, theme):
    """Add a progress bar at the bottom."""
    bar_w = 11.8
    _rect(slide, 0.78, 6.85, bar_w, 0.04, theme["border"])
    _rect(slide, 0.78, 6.85, bar_w * index / max(total, 1), 0.04, theme["primary"])


def _page_number(slide, index, total, theme):
    """Add page number in corner."""
    _txt(slide, f"{index:02d} / {total:02d}", 11.2, 6.8, 1.3, 0.3, 10, theme["muted"], True, PP_ALIGN.RIGHT)


# ── Cover slides (3 styles) ─────────────────────────────────────────

def _add_cover(prs, title, subtitle, author, theme, style="geometric"):
    if style == "wave":
        _cover_wave(prs, title, subtitle, author, theme)
    elif style == "minimal":
        _cover_minimal(prs, title, subtitle, author, theme)
    else:
        _cover_geometric(prs, title, subtitle, author, theme)


def _cover_geometric(prs, title, subtitle, author, theme):
    slide = _blank(prs, theme, dark=True)
    # Large decorative shapes
    _circle(slide, 8.5, -1.5, 5.5, theme["primary"])
    _circle(slide, 10.0, 4.0, 3.0, theme["secondary"])
    _circle(slide, 9.0, 2.5, 1.5, theme["accent"])
    _diamond(slide, 7.5, 0.5, 1.8, theme["gradient"][1])
    # Small scattered dots
    _circle(slide, 7.0, 1.0, 0.25, theme["text_light"])
    _circle(slide, 8.0, 5.5, 0.18, theme["text_light"])
    _circle(slide, 11.5, 3.5, 0.15, theme["text_light"])
    # Brand badge
    _rect(slide, 0.8, 0.7, 2.0, 0.45, theme["primary"], radius=True)
    _txt(slide, f"🚀 {author}", 0.95, 0.82, 1.7, 0.2, 11, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)
    # Title block
    _txt(slide, title, 0.85, 2.0, 7.0, 1.6, 40, RGBColor(255, 255, 255), True)
    _rect(slide, 0.85, 3.85, 4.5, 0.06, theme["accent"])
    _txt(slide, subtitle, 0.85, 4.1, 6.5, 0.6, 17, theme["text_light"])
    # Footer info
    _txt(slide, "AI 生成内容 · 专业排版 · 可编辑 PPTX", 0.85, 5.8, 5.5, 0.3, 11, theme["text_light"], italic=True)
    _decor_dots(slide, 0.85, 6.3, theme)


def _cover_wave(prs, title, subtitle, author, theme):
    slide = _blank(prs, theme, dark=True)
    # Wave-like layered rectangles (simulating curved layers)
    _rect(slide, 0, 4.8, 13.333, 2.7, theme["primary"])
    _rect(slide, 0, 5.3, 13.333, 2.2, theme["secondary"])
    _rect(slide, 0, 5.8, 13.333, 1.7, theme["accent"])
    # Floating circles
    _circle(slide, 10.5, 0.5, 2.5, theme["primary"])
    _circle(slide, 11.5, 2.0, 1.2, theme["gradient"][1])
    _circle(slide, 0.5, 0.5, 0.8, theme["secondary"])
    # Content
    _rect(slide, 0.9, 0.7, 2.2, 0.45, theme["accent"], radius=True)
    _txt(slide, f"📋 {author}", 1.05, 0.82, 2.0, 0.2, 11, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)
    _txt(slide, title, 0.9, 1.8, 8.5, 1.6, 42, RGBColor(255, 255, 255), True)
    _txt(slide, subtitle, 0.95, 3.5, 7.5, 0.6, 16, theme["text_light"])
    _txt(slide, "Powered by Lingxi AI", 0.95, 4.2, 4.0, 0.3, 11, theme["text_light"], italic=True)


def _cover_minimal(prs, title, subtitle, author, theme):
    slide = _blank(prs, theme, dark=False)
    # Left color bar
    _rect(slide, 0, 0, 0.35, 7.5, theme["primary"])
    # Subtle corner decoration
    _circle(slide, 10.5, -1.0, 4.0, theme["primary_light"])
    _diamond(slide, 11.5, 5.5, 1.5, theme["border"])
    # Content centered
    _txt(slide, author, 1.2, 1.5, 4.0, 0.4, 13, theme["muted"], True)
    _decor_line(slide, 1.2, 2.0, 2.5, theme)
    _txt(slide, title, 1.2, 2.5, 9.0, 1.6, 38, theme["dark"], True)
    _txt(slide, subtitle, 1.2, 4.3, 8.0, 0.6, 17, theme["muted"])
    # Bottom decorative elements
    _rect(slide, 1.2, 6.2, 11.0, 0.04, theme["border"])
    _txt(slide, "AI 自动生成 · 可编辑", 1.2, 6.4, 4.0, 0.3, 10, theme["muted"], italic=True)


# ── Table of contents ────────────────────────────────────────────────

def _add_toc_slide(prs, slides_data, theme):
    slide = _blank(prs, theme)
    _decor_corner_shapes(slide, theme)
    _rect(slide, 0.75, 0.5, 2.3, 0.45, theme["primary"], radius=True)
    _txt(slide, "📑 目录 CONTENTS", 0.92, 0.62, 2.0, 0.2, 11, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)
    _decor_line(slide, 0.75, 1.15, 11.5, theme, thin=False)

    cols = 2
    per_col = (len(slides_data) + 1) // cols
    for idx, item in enumerate(slides_data):
        col = idx // per_col
        row = idx % per_col
        x = 0.85 + col * 6.2
        y = 1.6 + row * 0.95
        # Number circle
        _circle(slide, x, y + 0.06, 0.48, theme["primary"] if col == 0 else theme["secondary"])
        _txt(slide, f"{idx + 1:02d}", x + 0.12, y + 0.18, 0.24, 0.14, 10, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)
        # Title
        kicker = item.get("kicker", "")
        title = item.get("title", f"第 {idx + 1} 章节")
        _txt(slide, title, x + 0.65, y + 0.05, 4.8, 0.3, 14, theme["dark"], True)
        if kicker:
            _txt(slide, kicker, x + 0.65, y + 0.4, 2.5, 0.2, 10, theme["muted"])
        # Separator line
        _rect(slide, x + 0.65, y + 0.7, 4.5, 0.015, theme["border"])


# ── Section divider ──────────────────────────────────────────────────

def _add_section_divider(prs, item, index, total, theme):
    slide = _blank(prs, theme, dark=True)
    _circle(slide, -1.0, -1.0, 4.0, theme["primary"])
    _circle(slide, 10.0, 4.5, 3.5, theme["secondary"])
    _diamond(slide, 5.5, 0.5, 1.5, theme["gradient"][1])
    kicker = item.get("kicker") or f"PART {index:02d}"
    _rect(slide, 5.2, 2.5, 2.8, 0.4, theme["accent"], radius=True)
    _txt(slide, kicker, 5.4, 2.58, 2.4, 0.2, 11, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)
    _txt(slide, item.get("title", ""), 2.5, 3.2, 8.3, 0.8, 32, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)
    _rect(slide, 5.5, 4.2, 2.3, 0.05, theme["accent"])


# ── Slide header (shared) ───────────────────────────────────────────

def _slide_header(slide, item, index, total, theme):
    kicker = item.get("kicker") or f"PART {index:02d}"
    _rect(slide, 0.75, 0.5, 1.7, 0.38, theme["primary"], radius=True)
    _txt(slide, kicker[:18], 0.88, 0.59, 1.4, 0.2, 9, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)
    _txt(slide, item.get("title", ""), 0.75, 1.05, 8.5, 0.55, 24, theme["dark"], True)
    _page_number(slide, index, total, theme)
    # Top-right accent line
    _rect(slide, 10.2, 1.1, 2.3, 0.06, theme["accent"])
    # Top-right dot
    _circle(slide, 12.2, 0.5, 0.35, theme["primary_light"])
    # Note text if provided
    note = item.get("note", "")
    if note:
        _txt(slide, note, 0.78, 6.35, 10.0, 0.3, 10, theme["muted"], italic=True)
    _progress_bar(slide, index, total, theme)


# ── Layout: Cards (2x2 grid with icons) ─────────────────────────────

def _add_cards_slide(prs, item, index, total, theme):
    slide = _blank(prs, theme)
    _decor_corner_shapes(slide, theme)
    _slide_header(slide, item, index, total, theme)
    bullets = item.get("bullets", [])[:4]
    cw, ch = 5.55, 1.55
    positions = [(0.78, 1.95), (6.7, 1.95), (0.78, 3.85), (6.7, 3.85)]
    icons = ["💡", "📊", "🎯", "⚡"]
    for idx, bullet in enumerate(bullets):
        left, top = positions[idx]
        # Card background with left color bar
        _rect(slide, left, top, cw, ch, theme["card"], theme["border"], radius=True)
        _rect(slide, left + 0.05, top + 0.15, 0.08, ch - 0.3, theme["primary"] if idx % 2 == 0 else theme["secondary"])
        # Icon circle
        _circle(slide, left + 0.3, top + 0.35, 0.7, theme["card_alt"])
        _txt(slide, icons[idx % 4], left + 0.42, top + 0.5, 0.45, 0.25, 16, theme["primary"], align=PP_ALIGN.CENTER)
        # Number badge
        _circle(slide, left + cw - 0.55, top + 0.15, 0.38, theme["primary"])
        _txt(slide, str(idx + 1), left + cw - 0.42, top + 0.24, 0.12, 0.1, 9, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)
        # Text content
        _txt(slide, bullet, left + 1.2, top + 0.3, cw - 1.8, 0.9, 14, theme["dark"], True)


# ── Layout: Timeline (horizontal) ───────────────────────────────────

def _add_timeline_slide(prs, item, index, total, theme):
    slide = _blank(prs, theme)
    _slide_header(slide, item, index, total, theme)
    bullets = item.get("bullets", [])[:5]
    n = len(bullets)
    # Timeline bar
    bar_y = 3.3
    _rect(slide, 0.95, bar_y, 11.4, 0.08, theme["primary_light"])
    gap = 10.5 / max(n - 1, 1)
    for idx, bullet in enumerate(bullets):
        x = 1.0 + idx * gap
        # Node circle with ring
        _circle(slide, x + 0.05, bar_y - 0.25, 0.62, theme["card"], theme["primary"])
        _circle(slide, x + 0.15, bar_y - 0.15, 0.42, theme["primary"])
        _txt(slide, str(idx + 1), x + 0.27, bar_y - 0.03, 0.18, 0.14, 10, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)
        # Card below
        card_x = x - 0.35
        _rect(slide, card_x, 3.85, 2.2, 1.5, theme["card"], theme["border"], radius=True)
        _rect(slide, card_x + 0.08, 3.95, 2.04, 0.32, theme["primary"] if idx % 2 == 0 else theme["secondary"], radius=True)
        _txt(slide, f"Step {idx + 1}", card_x + 0.2, 3.98, 1.7, 0.2, 9, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)
        _txt(slide, bullet, card_x + 0.15, 4.4, 1.9, 0.7, 11, theme["dark"], True, PP_ALIGN.CENTER)
    # Decorative elements
    _diamond(slide, 12.0, 1.8, 0.6, theme["primary_light"])


# ── Layout: Process (left-to-right arrow flow) ──────────────────────

def _add_process_slide(prs, item, index, total, theme):
    slide = _blank(prs, theme)
    _slide_header(slide, item, index, total, theme)
    bullets = item.get("bullets", [])[:5]
    n = len(bullets)
    step_w = 2.0
    gap = (11.5 - n * step_w) / max(n, 1)
    start_x = 0.9
    for idx, bullet in enumerate(bullets):
        x = start_x + idx * (step_w + gap)
        y = 2.4
        # Step box
        color = theme["primary"] if idx % 2 == 0 else theme["secondary"]
        _rect(slide, x, y, step_w, 3.2, theme["card"], theme["border"], radius=True)
        # Top colored header bar
        _rect(slide, x + 0.06, y + 0.08, step_w - 0.12, 0.55, color, radius=True)
        _txt(slide, f"0{idx + 1}", x + 0.25, y + 0.17, 0.5, 0.25, 14, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)
        # Icon
        process_icons = ["📋", "⚙️", "🔍", "✅", "🚀"]
        _txt(slide, process_icons[idx % 5], x + 0.65, y + 0.85, 0.7, 0.5, 22, theme["primary"], align=PP_ALIGN.CENTER)
        # Content
        _txt(slide, bullet, x + 0.15, y + 1.5, step_w - 0.3, 1.4, 12, theme["dark"], True, PP_ALIGN.CENTER)
        # Arrow between steps
        if idx < n - 1:
            arrow_x = x + step_w + gap * 0.2
            _arrow_right(slide, arrow_x, 3.6, gap * 0.6, 0.5, theme["primary_light"])


# ── Layout: Metrics (data visualization style) ──────────────────────

def _add_metrics_slide(prs, item, index, total, theme):
    slide = _blank(prs, theme)
    _slide_header(slide, item, index, total, theme)
    bullets = item.get("bullets", [])[:4]
    n = len(bullets)
    card_w = 2.65
    gap = (11.5 - n * card_w) / (n + 1)
    for idx, bullet in enumerate(bullets):
        x = 0.9 + gap + idx * (card_w + gap)
        y = 2.2
        # Card with top border color
        _rect(slide, x, y, card_w, 3.5, theme["card"], theme["border"], radius=True)
        _rect(slide, x + 0.05, y + 0.05, card_w - 0.1, 0.12, theme["primary"] if idx % 2 == 0 else theme["secondary"])
        # Large number
        _txt(slide, f"0{idx + 1}", x + 0.3, y + 0.5, 1.0, 0.6, 32, theme["primary"], True)
        # Hexagon decoration
        _hexagon(slide, x + card_w - 0.8, y + 0.4, 0.55, 0.55, theme["primary_light"])
        # Divider
        _rect(slide, x + 0.3, y + 1.3, card_w - 0.6, 0.025, theme["border"])
        # Content
        _txt(slide, bullet, x + 0.25, y + 1.6, card_w - 0.5, 1.5, 13, theme["dark"], True)
        # Bottom bar (simulating a metric bar)
        bar_heights = [0.6, 0.8, 0.5, 0.9]
        bh = bar_heights[idx % 4]
        _rect(slide, x + 0.4, y + 3.5 - bh - 0.2, 0.3, bh, theme["primary"])
        _rect(slide, x + 0.8, y + 3.5 - bh * 0.7 - 0.2, 0.3, bh * 0.7, theme["secondary"])
        _rect(slide, x + 1.2, y + 3.5 - bh * 0.5 - 0.2, 0.3, bh * 0.5, theme["primary_light"])


# ── Layout: Comparison (two-column) ─────────────────────────────────

def _add_comparison_slide(prs, item, index, total, theme):
    slide = _blank(prs, theme)
    _slide_header(slide, item, index, total, theme)
    bullets = item.get("bullets", [])[:4]
    # Left column
    _rect(slide, 0.78, 1.9, 5.7, 4.3, theme["card"], theme["border"], radius=True)
    _rect(slide, 0.83, 1.95, 5.6, 0.5, theme["primary"], radius=True)
    _txt(slide, "✦ 核心要点", 1.1, 2.02, 2.5, 0.3, 13, RGBColor(255, 255, 255), True)
    for idx, bullet in enumerate(bullets[:2]):
        y = 2.7 + idx * 1.4
        _circle(slide, 1.1, y + 0.1, 0.5, theme["primary_light"])
        _txt(slide, str(idx + 1), 1.28, y + 0.22, 0.15, 0.12, 10, theme["primary"], True, PP_ALIGN.CENTER)
        _txt(slide, bullet, 1.8, y + 0.05, 4.3, 0.9, 14, theme["dark"], True)
    # Right column
    _rect(slide, 6.85, 1.9, 5.7, 4.3, theme["card"], theme["border"], radius=True)
    _rect(slide, 6.9, 1.95, 5.6, 0.5, theme["secondary"], radius=True)
    _txt(slide, "✦ 延伸分析", 7.15, 2.02, 2.5, 0.3, 13, RGBColor(255, 255, 255), True)
    for idx, bullet in enumerate(bullets[2:4]):
        y = 2.7 + idx * 1.4
        _circle(slide, 7.15, y + 0.1, 0.5, theme["primary_light"])
        _txt(slide, str(idx + 3), 7.33, y + 0.22, 0.15, 0.12, 10, theme["secondary"], True, PP_ALIGN.CENTER)
        _txt(slide, bullet, 7.85, y + 0.05, 4.3, 0.9, 14, theme["dark"], True)
    _progress_bar(slide, index, total, theme)


# ── Layout: Icon Grid (3x2 or 2x3) ─────────────────────────────────

def _add_icon_grid_slide(prs, item, index, total, theme):
    slide = _blank(prs, theme)
    _slide_header(slide, item, index, total, theme)
    bullets = item.get("bullets", [])[:6]
    grid_icons = ["🎯", "📈", "💡", "🔧", "🌟", "🚀"]
    cols = 3
    rows = (len(bullets) + cols - 1) // cols
    cell_w = 3.6
    cell_h = 1.8
    start_x = 0.85
    start_y = 2.0
    gap_x = 0.3
    gap_y = 0.3
    for idx, bullet in enumerate(bullets):
        col = idx % cols
        row = idx // cols
        x = start_x + col * (cell_w + gap_x)
        y = start_y + row * (cell_h + gap_y)
        _rect(slide, x, y, cell_w, cell_h, theme["card"], theme["border"], radius=True)
        # Icon circle
        _circle(slide, x + 0.2, y + 0.2, 0.65, theme["primary"] if idx % 2 == 0 else theme["secondary"])
        _txt(slide, grid_icons[idx % 6], x + 0.34, y + 0.33, 0.35, 0.2, 14, RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
        # Text
        _txt(slide, bullet, x + 1.05, y + 0.25, cell_w - 1.3, cell_h - 0.5, 13, theme["dark"], True)


# ── Closing slide ────────────────────────────────────────────────────

def _add_closing(prs, title, author, theme):
    slide = _blank(prs, theme, dark=True)
    # Decorative circles
    _circle(slide, -0.5, -0.5, 3.5, theme["primary"])
    _circle(slide, 0.5, 5.0, 2.0, theme["secondary"])
    _circle(slide, 9.5, 3.5, 4.0, theme["gradient"][1])
    _circle(slide, 11.0, 0.5, 1.5, theme["accent"])
    _diamond(slide, 6.0, 0.5, 1.0, theme["primary_light"])
    # Thank you text
    _txt(slide, "THANK YOU", 2.5, 2.0, 8.3, 0.8, 42, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)
    _txt(slide, "谢谢观看", 2.5, 2.9, 8.3, 0.6, 22, theme["text_light"], False, PP_ALIGN.CENTER)
    _rect(slide, 5.5, 3.7, 2.3, 0.05, theme["accent"])
    _txt(slide, f"Powered by {author}", 3.5, 4.0, 6.3, 0.35, 13, theme["text_light"], italic=True, align=PP_ALIGN.CENTER)
    _txt(slide, title, 3.5, 4.5, 6.3, 0.35, 12, theme["text_light"], align=PP_ALIGN.CENTER)
    # Bottom decorative dots
    for i in range(8):
        _circle(slide, 4.5 + i * 0.55, 5.5, 0.12, theme["text_light"])


# =====================================================================
#  PDF generator
# =====================================================================

def generate_pdf(data: dict[str, Any]) -> Path:
    title_text = data.get("title", "文档")
    sections = data.get("sections", [])
    filename = f"{_uid()}_{_safe(title_text)}.pdf"
    path = OUTPUT_DIR / filename

    font_name = _register_cjk_font()

    doc = SimpleDocTemplate(
        str(path),
        pagesize=A4,
        leftMargin=25 * mm,
        rightMargin=25 * mm,
        topMargin=25 * mm,
        bottomMargin=25 * mm,
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "DocTitle",
        parent=styles["Title"],
        fontName=font_name,
        fontSize=22,
        leading=28,
        spaceAfter=20,
        alignment=1,  # center
    )
    heading_style = ParagraphStyle(
        "DocHeading",
        parent=styles["Heading2"],
        fontName=font_name,
        fontSize=15,
        leading=20,
        spaceBefore=16,
        spaceAfter=8,
        textColor=HexColor("#1e40af"),
    )
    body_style = ParagraphStyle(
        "DocBody",
        parent=styles["Normal"],
        fontName=font_name,
        fontSize=11,
        leading=18,
        spaceAfter=8,
    )
    footer_style = ParagraphStyle(
        "DocFooter",
        parent=styles["Normal"],
        fontName=font_name,
        fontSize=9,
        leading=12,
        alignment=2,  # right
        textColor=HexColor("#64748b"),
    )

    story: list[Any] = []
    story.append(Paragraph(title_text, title_style))
    story.append(Spacer(1, 12))

    for sec in sections:
        story.append(Paragraph(sec.get("heading", ""), heading_style))
        for para in sec.get("paragraphs", []):
            story.append(Paragraph(para, body_style))

    story.append(Spacer(1, 30))
    story.append(Paragraph("—— 由灵犀 AI 自动生成", footer_style))

    doc.build(story)
    return path


# =====================================================================
#  Excel generator
# =====================================================================

def generate_xlsx(data: dict[str, Any]) -> Path:
    title_text = data.get("title", "表格")
    headers = data.get("headers", [])
    rows = data.get("rows", [])
    filename = f"{_uid()}_{_safe(title_text)}.xlsx"
    path = OUTPUT_DIR / filename

    wb = Workbook()
    ws = wb.active
    ws.title = title_text[:31]  # Excel sheet name max 31 chars

    # ── Title row ──
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=max(len(headers), 1))
    title_cell = ws.cell(row=1, column=1, value=title_text)
    title_cell.font = Font(size=16, bold=True, color="1e40af")
    title_cell.alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 35

    # ── Header row ──
    header_fill = PatternFill(start_color="1e40af", end_color="1e40af", fill_type="solid")
    header_font = Font(size=12, bold=True, color="ffffff")
    header_align = Alignment(horizontal="center", vertical="center")
    thin_border = Border(
        left=Side(style="thin"),
        right=Side(style="thin"),
        top=Side(style="thin"),
        bottom=Side(style="thin"),
    )
    for col_idx, header in enumerate(headers, 1):
        cell = ws.cell(row=2, column=col_idx, value=header)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = header_align
        cell.border = thin_border
        ws.column_dimensions[cell.column_letter].width = max(len(str(header)) * 2 + 6, 14)
    ws.row_dimensions[2].height = 28

    # ── Data rows ──
    data_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
    for row_idx, row in enumerate(rows, 3):
        for col_idx, value in enumerate(row, 1):
            cell = ws.cell(row=row_idx, column=col_idx, value=value)
            cell.alignment = data_align
            cell.border = thin_border
        ws.row_dimensions[row_idx].height = 22

    wb.save(str(path))
    return path


# =====================================================================
#  Dispatcher
# =====================================================================

def generate_document(doc_type: str, data: dict[str, Any]) -> Path:
    """Generate a real document file and return its path."""
    if doc_type == "pptx":
        return generate_pptx(data)
    if doc_type == "pdf":
        return generate_pdf(data)
    if doc_type == "xlsx":
        return generate_xlsx(data)
    raise ValueError(f"Unsupported doc_type: {doc_type}")


def _safe(text: str) -> str:
    """Sanitize filename component."""
    s = re.sub(r"[^\w\u4e00-\u9fff-]", "_", text)
    return s[:30]
